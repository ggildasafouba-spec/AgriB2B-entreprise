#!/usr/bin/env python3
"""
Script pour générer automatiquement la présentation PowerPoint d'AgroMarket
Nécessite : pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_agromarket_presentation():
    prs = Presentation()

    # Slide 1: Titre
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "AgroMarket"
    subtitle.text = "Révolutionner l'Agriculture Africaine\nLe marché agricole au bout des doigts"

    # Slide 2: Problèmes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Les Défis de l'Agriculture Africaine"

    content = slide.placeholders[1]
    content.text = """• Accès limité aux marchés pour les agriculteurs
• Intermédiaires coûteux réduisant les revenus
• Paiements incertains et délais importants
• Manque de transparence dans la chaîne
• Volatilité des prix non maîtrisée"""

    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Notre Solution : AgroMarket"

    content = slide.placeholders[1]
    content.text = """Plateforme digitale B2B connectant directement
producteurs et acheteurs africains

• Élimination des intermédiaires (+30-50% revenus)
• Paiements mobiles intégrés (MTN MoMo, Orange Money)
• Système d'escrow sécurisé
• Gestion intelligente des stocks et commandes"""

    # Slide 4: Fonctionnalités
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Fonctionnalités Clés"

    content = slide.placeholders[1]
    content.text = """• Gestion des Produits avec catalogue digital
• Suivi des Stocks en temps réel
• Système de Commandes intuitif
• Paiements Sécurisés avec escrow
• KYC Intégré pour la confiance
• Tableau de Bord Analytics
• Notifications Temps Réel
• Support Multi-Rôles"""

    # Slide 5: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Architecture Technique Robuste"

    content = slide.placeholders[1]
    content.text = """Backend:
• NestJS + TypeScript
• Prisma ORM + PostgreSQL
• Redis pour le cache
• JWT Authentication

Frontend:
• Next.js 14 + React
• Tailwind CSS
• API REST

Infrastructure:
• Docker + Docker Compose
• Intégrations Mobile Money"""

    # Slide 6: Impact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Impact Économique Prévu"

    content = slide.placeholders[1]
    content.text = """• +30-50% de revenus pour les producteurs
• -20-30% sur les prix grâce aux intermédiaires
• Paiements instantanés vs délais de 90 jours
• Création d'emplois dans la tech agricole
• Développement économique rural"""

    # Slide 7: Marché et Stratégie
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Marché Cible et Stratégie"

    content = slide.placeholders[1]
    content.text = """Segments cibles:
• Agriculteurs africains (petits/moyens)
• Acheteurs B2B (commerçants, entreprises)

Stratégie Go-to-Market:
• Phase 1: Pilote Sénégal (Q2 2026)
• Phase 2: Expansion régionale (Q4 2026)
• Phase 3: Scale national (2027)"""

    # Slide 8: Modèle Économique
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Modèle Économique Durable"

    content = slide.placeholders[1]
    content.text = """Revenus:
• Commission 5% sur transactions
• Services premium (analytics, support)
• Partenariats bancaires

Coûts:
• Infrastructure et maintenance
• Équipe technique et support
• Conformité et sécurité

Projection: Rentable à 24 mois"""

    # Slide 9: Équipe
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Équipe et Expertise"

    content = slide.placeholders[1]
    content.text = """• Développeur Full-Stack (React, Node.js)
• Spécialiste Agriculture africaine
• Expert Paiements Mobiles (Fintech)
• Designer UX/UI professionnel

Compétences clés:
• Développement web moderne
• Intégration API paiements
• Gestion bases de données
• Déploiement cloud"""

    # Slide 10: Démo
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Démonstration Live"

    content = slide.placeholders[1]
    content.text = """Points à démontrer:
1. Page d'accueil avec branding
2. Inscription et KYC
3. Dashboard vendeur (produits/stocks)
4. Interface acheteur (commandes)
5. Paiement mobile money
6. Analytics et rapports"""

    # Slide 11: Concurrents
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Analyse Concurrentielle"

    content = slide.placeholders[1]
    content.text = """Concurrents:
• Solutions locales limitées
• Marketplaces internationales inadaptées
• Systèmes traditionnels (marchés physiques)

Avantages concurrentiels:
• Spécialisation africaine
• Paiements mobiles natifs
• Interface multilingue
• Prix abordable"""

    # Slide 12: Risques
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Risques et Mitigation"

    content = slide.placeholders[1]
    content.text = """Risques identifiés:
• Sécurité paiements → Chiffrement SSL, conformité
• Adoption lente → Marketing local, démos terrain
• Réglementation → Conformité lois locales
• Concurrence → Innovation continue, partenariats"""

    # Slide 13: Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Feuille de Route 2026-2028"

    content = slide.placeholders[1]
    content.text = """Q2 2026: Lancement pilote Sénégal
Q4 2026: Expansion Côte d'Ivoire/Cameroun
2027: Scale national (5 pays)
2028: Expansion continentale + fonctionnalités avancées

Objectif: 100K utilisateurs actifs d'ici 2028"""

    # Slide 14: Call to Action
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Rejoignez l'Aventure AgroMarket"

    content = slide.placeholders[1]
    content.text = """Opportunités:
• Investisseurs: Transformer l'agriculture africaine
• Partenaires: Développer l'écosystème
• Utilisateurs: Bénéficier d'une plateforme innovante

Contact: contact@agromarket.africa
Site: www.agromarket.africa"""

    # Slide 15: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Questions & Réponses"

    content = slide.placeholders[1]
    content.text = """Espace pour vos questions

Points préparés:
• Aspects techniques détaillés
• Modèle économique approfondi
• Stratégie d'adoption
• Analyse concurrentielle
• Prochaines étapes"""

    # Sauvegarder
    prs.save('AgroMarket_Presentation.pptx')
    print("✅ Présentation PowerPoint créée: AgroMarket_Presentation.pptx")

if __name__ == "__main__":
    create_agromarket_presentation()